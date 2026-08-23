"""Build VERA-Complete.pptx — widescreen 16:9 PowerPoint for managers."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
SHOTS = ROOT / "screenshots"
OUT = ROOT / "VERA-Complete.pptx"
OUT_FALLBACKS = [
    ROOT / "VERA-Complete-updated.pptx",
    ROOT / "VERA-Complete-latest.pptx",
    ROOT / "VERA-Complete-v3.pptx",
    ROOT / "VERA-Complete-v4.pptx",
    ROOT / "VERA-Complete-manager.pptx",
    ROOT / "VERA-Complete-final.pptx",
    ROOT / "VERA-Complete-fixed.pptx",
]

# 16:9 widescreen
W, H = Inches(13.333), Inches(7.5)

NAVY = RGBColor(0x0A, 0x24, 0x30)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
TEAL2 = RGBColor(0x14, 0xB8, 0xA6)
CORAL = RGBColor(0xE8, 0x5D, 0x04)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF7, 0xF4, 0xEF)
MUTED = RGBColor(0x5C, 0x71, 0x7A)
DARK = RGBColor(0x06, 0x15, 0x1C)
LIGHT_TEAL = RGBColor(0xEC, 0xFD, 0xF5)
LIGHT_CORAL = RGBColor(0xFF, 0xF7, 0xED)


def set_run(run, size=18, bold=False, color=NAVY, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def fill_solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def send_back(slide, shape):
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_bg(slide, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    fill_solid(shape, color)
    send_back(slide, shape)


def add_light_chrome(slide, accent=TEAL):
    """Polished light slide: cream bg, left accent, soft top band."""
    add_bg(slide, CREAM)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, Inches(0.08))
    fill_solid(band, accent)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.14), H)
    fill_solid(bar, accent)
    footer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(7.15), W, Inches(0.35)
    )
    fill_solid(footer, RGBColor(0xE8, 0xEE, 0xF0))
    textbox(
        slide,
        Inches(0.45),
        Inches(7.18),
        Inches(8),
        Inches(0.28),
        "VERA  ·  Verified Evidence, Reliable Agents",
        size=10,
        color=MUTED,
    )


def add_accent_bar(slide, color=TEAL):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.12), H)
    fill_solid(bar, color)


def textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    """Add text; split on newlines into real paragraphs (avoids clipped mid-line text)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = None
    except Exception:
        pass
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = str(text).split("\n") if text is not None else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        p.text = line
        for run in p.runs:
            set_run(run, size=size, bold=bold, color=color)
    return box


def flow_row(slide, items, *, top, height, left=Inches(0.4), right_margin=Inches(0.4), gap=Inches(0.3), detail_size=12):
    """Lay out flow_box items in a row with real gaps and centered arrows (no overlap)."""
    n = len(items)
    if n == 0:
        return
    usable = W - left - right_margin
    box_w = (usable - gap * (n - 1)) / n
    x = left
    for i, item in enumerate(items):
        title, sub, color, detail = item[0], item[1], item[2], item[3] if len(item) > 3 else ""
        flow_box(slide, x, top, box_w, height, title, sub, color, detail, detail_size=detail_size)
        if i < n - 1:
            textbox(
                slide,
                x + box_w,
                top + height / 2 - Inches(0.18),
                gap,
                Inches(0.36),
                "→",
                size=14,
                bold=True,
                color=TEAL2,
                align=PP_ALIGN.CENTER,
            )
        x += box_w + gap


def card(slide, left, top, width, height, title, body, *, fill=WHITE, title_color=NAVY, body_size=12):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_solid(shape, fill)
    shape.line.color.rgb = RGBColor(0xD0, 0xDB, 0xE0)
    shape.adjustments[0] = 0.08
    # thin accent on card top
    tip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.06)
    )
    fill_solid(tip, title_color if title_color != NAVY else TEAL)
    textbox(
        slide,
        left + Inches(0.18),
        top + Inches(0.16),
        width - Inches(0.36),
        Inches(0.36),
        title,
        size=14,
        bold=True,
        color=title_color,
    )
    textbox(
        slide,
        left + Inches(0.18),
        top + Inches(0.52),
        width - Inches(0.36),
        height - Inches(0.65),
        body,
        size=body_size,
        color=MUTED,
    )


def visual_card(
    slide,
    left,
    top,
    width,
    height,
    title,
    body,
    image_name: str,
    *,
    fill=WHITE,
    title_color=NAVY,
):
    """Card with title + short body + product screenshot filling remaining space."""
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_solid(shape, fill)
    shape.line.color.rgb = RGBColor(0xD0, 0xDB, 0xE0)
    shape.adjustments[0] = 0.06
    tip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.07))
    fill_solid(tip, title_color if title_color != NAVY else TEAL)
    textbox(
        slide,
        left + Inches(0.2),
        top + Inches(0.18),
        width - Inches(0.4),
        Inches(0.35),
        title,
        size=16,
        bold=True,
        color=title_color,
    )
    textbox(
        slide,
        left + Inches(0.2),
        top + Inches(0.52),
        width - Inches(0.4),
        Inches(0.85),
        body,
        size=11,
        color=MUTED,
    )
    img_top = top + Inches(1.45)
    img_h = height - Inches(1.6)
    img_w = width - Inches(0.3)
    add_image_if(
        slide,
        image_name,
        left + Inches(0.15),
        img_top,
        img_w,
        max_height=img_h,
    )


def flow_box(slide, left, top, w, h, title, sub, color, detail: str = "", *, detail_size=11):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, w, h)
    fill_solid(shape, color)
    shape.adjustments[0] = 0.12
    textbox(slide, left + Inches(0.06), top + Inches(0.16), w - Inches(0.12), Inches(0.32), title, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(
        slide,
        left + Inches(0.08),
        top + Inches(0.46),
        w - Inches(0.16),
        Inches(0.28),
        sub,
        size=10,
        bold=True,
        color=RGBColor(0xE8, 0xF5, 0xF4),
        align=PP_ALIGN.CENTER,
    )
    if detail:
        textbox(
            slide,
            left + Inches(0.1),
            top + Inches(0.82),
            w - Inches(0.2),
            h - Inches(0.95),
            detail,
            size=detail_size,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )


def arrow(slide, left, top):
    textbox(slide, left, top, Inches(0.35), Inches(0.4), "→", size=22, bold=True, color=TEAL2, align=PP_ALIGN.CENTER)


def add_image_if(slide, name: str, left, top, width, height=None, *, max_height=None):
    """Place image; keep aspect ratio when height is omitted or max_height is set."""
    path = SHOTS / name
    if path.exists():
        if height is None:
            pic = slide.shapes.add_picture(str(path), left, top, width=width)
            if max_height is not None and pic.height > max_height:
                ratio = max_height / pic.height
                pic.height = max_height
                pic.width = int(pic.width * ratio)
            return True
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        return True
    h = height or Inches(4)
    missing = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, h)
    fill_solid(missing, CREAM)
    textbox(slide, left, top + h / 2 - Inches(0.2), width, Inches(0.4), f"Missing {name}", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    return False


def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    # Atmosphere — large soft shapes (no purple glow clutter)
    c1 = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.8), Inches(-2.2), Inches(6.5), Inches(6.5))
    fill_solid(c1, RGBColor(0x0E, 0x3A, 0x3A))
    c2 = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-2.2), Inches(4.2), Inches(5.2), Inches(5.2))
    fill_solid(c2, RGBColor(0x3A, 0x1C, 0x08))
    # thin coral rail
    rail = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.16), H)
    fill_solid(rail, CORAL)

    textbox(s, Inches(0.85), Inches(1.15), Inches(11), Inches(0.35), "VERIFIED EVIDENCE  ·  RELIABLE AGENTS", size=12, bold=True, color=TEAL2)
    textbox(s, Inches(0.8), Inches(1.7), Inches(12), Inches(1.1), "VERA", size=72, bold=True, color=WHITE)
    textbox(
        s,
        Inches(0.85),
        Inches(3.0),
        Inches(11.5),
        Inches(1.5),
        "Stop shipping chatbots that sound right.\nShip agents that can prove they’re right.",
        size=28,
        bold=True,
        color=WHITE,
    )
    textbox(
        s,
        Inches(0.85),
        Inches(4.8),
        Inches(11.2),
        Inches(0.9),
        "Upload your knowledge. VERA becomes the careful librarian — corkboard, sticky notes, and receipts —\nthen embeds as a specialist your customers can trust on any website.",
        size=16,
        color=RGBColor(0xB8, 0xCB, 0xD1),
    )
    textbox(
        s,
        Inches(0.85),
        Inches(6.35),
        Inches(11.5),
        Inches(0.4),
        "Knowledge Graph   ·   Hybrid Ask   ·   Trust Trail   ·   Prove it   ·   Embed anywhere",
        size=14,
        bold=True,
        color=CORAL,
    )


def section(prs, num: str, title: str, subtitle: str = ""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CREAM)
    add_accent_bar(s, TEAL)
    textbox(s, Inches(0.7), Inches(2.4), Inches(11), Inches(0.4), num, size=14, bold=True, color=CORAL)
    textbox(s, Inches(0.7), Inches(2.9), Inches(11), Inches(1.2), title, size=36, bold=True, color=NAVY)
    if subtitle:
        textbox(s, Inches(0.7), Inches(4.3), Inches(10), Inches(1), subtitle, size=18, color=MUTED)


def content_title(s, kicker: str, title: str):
    textbox(s, Inches(0.5), Inches(0.28), Inches(12), Inches(0.28), kicker, size=11, bold=True, color=TEAL)
    textbox(s, Inches(0.5), Inches(0.55), Inches(12.2), Inches(0.7), title, size=26, bold=True, color=NAVY)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    title_slide(prs)

    # What is VERA — visual pillars with real Studio shots
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_chrome(s)
    content_title(s, "01  ·  WHAT VERA IS", "Upload the library. Ship the librarian.")
    visual_card(
        s,
        Inches(0.4),
        Inches(1.35),
        Inches(4.05),
        Inches(5.35),
        "1 · Ingest",
        "Bring data in — files, websites, SharePoint. Cleaned and woven into the agent’s private world.",
        "thumbs/thumb-ingest.png",
        title_color=TEAL,
    )
    visual_card(
        s,
        Inches(4.65),
        Inches(1.35),
        Inches(4.05),
        Inches(5.35),
        "2 · Intelligence",
        "Knowledge graph + vector memory — grounded answers with a Trust Trail.",
        "thumbs/thumb-intelligence.png",
        title_color=TEAL,
        fill=LIGHT_TEAL,
    )
    visual_card(
        s,
        Inches(8.9),
        Inches(1.35),
        Inches(4.05),
        Inches(5.35),
        "3 · Ship",
        "Publish & embed a trusted specialist on any website — same brain, daily use.",
        "thumbs/thumb-ship.png",
        title_color=CORAL,
        fill=LIGHT_CORAL,
    )

    # Problem
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_chrome(s, CORAL)
    content_title(s, "02  ·  THE PROBLEM", "AI is fluent. Trust is scarce.")
    card(s, Inches(0.55), Inches(1.55), Inches(5.9), Inches(4.8), "Typical chatbot / vector RAG — “Trust me”", "• Finds text that sounds similar\n• Misses multi-hop relationships\n• Smooths over conflicts\n• Hard for managers & clients to audit\n• Similarity score ≠ proof", fill=LIGHT_CORAL, title_color=CORAL)
    card(s, Inches(6.8), Inches(1.55), Inches(5.9), Inches(4.8), "VERA — “Here’s the evidence”", "• Builds a knowledge graph from your data\n• Keeps vector memory for exact wording\n• Surfaces conflicts & unsupported claims\n• Trust Trail + Prove it on demand\n• Clarify / refuse when evidence is thin", fill=LIGHT_TEAL, title_color=TEAL)

    # Magic flow — short lines + real gaps so text never clips under arrows
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    textbox(s, Inches(0.5), Inches(0.28), Inches(12), Inches(0.28), "03  ·  THE MAGIC", size=11, bold=True, color=TEAL2)
    textbox(s, Inches(0.5), Inches(0.58), Inches(12.2), Inches(0.55), "Your books in. A librarian out.", size=28, bold=True, color=WHITE)
    textbox(
        s,
        Inches(0.5),
        Inches(1.15),
        Inches(12),
        Inches(0.35),
        "One pipeline — raw knowledge → proven specialist → live on any site",
        size=15,
        color=RGBColor(0xB8, 0xCB, 0xD1),
    )
    steps = [
        ("1. Connect", "Upload · Web · SP", TEAL, "Bring in files,\nsites, or SharePoint."),
        ("2. Clean", "CleanStack", TEAL2, "Strip noise so\nthe text is readable."),
        ("3. Weave", "Knowledge Graph", TEAL, "Build the corkboard\nof real relationships."),
        ("4. Embed", "Vector (Chroma)", TEAL2, "Store sticky-note\nquotes for recall."),
        ("5. Ask", "Trust Trail", CORAL, "Answer with proof\n— or refuse."),
        ("6. Publish", "Website widget", RGBColor(0xF4, 0x8C, 0x06), "Embed the same\nagent on any site."),
    ]
    flow_row(s, steps, top=Inches(1.7), height=Inches(3.35))
    textbox(
        s,
        Inches(0.5),
        Inches(5.3),
        Inches(12.2),
        Inches(1.3),
        "Same brain in Studio and on the client site — no second system.\n"
        "Don’t ask them to trust the model. Make the model show the map.",
        size=16,
        bold=True,
        color=WHITE,
    )

    # Two memories
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_chrome(s)
    content_title(s, "04  ·  TWO MEMORIES", "Corkboard + sticky notes")
    # left corkboard
    board = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.7), Inches(5.9), Inches(4.8))
    fill_solid(board, TEAL)
    board.adjustments[0] = 0.08
    textbox(s, Inches(0.85), Inches(1.95), Inches(5), Inches(0.3), "CORKBOARD", size=11, bold=True, color=TEAL2)
    textbox(s, Inches(0.85), Inches(2.3), Inches(5), Inches(0.5), "Knowledge Graph", size=24, bold=True, color=WHITE)
    n1 = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(1.4), Inches(3.4), Inches(1.6), Inches(1.6))
    fill_solid(n1, WHITE)
    textbox(s, Inches(1.4), Inches(4.0), Inches(1.6), Inches(0.4), "Cap’n Bill", size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    n2 = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(3.8), Inches(3.2), Inches(1.5), Inches(1.5))
    fill_solid(n2, RGBColor(0xCC, 0xFB, 0xF1))
    textbox(s, Inches(3.8), Inches(3.75), Inches(1.5), Inches(0.4), "Tin Woodman", size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    textbox(s, Inches(0.85), Inches(5.7), Inches(5.2), Inches(0.5), "Who connects to whom — with proof on the strings", size=13, color=WHITE)
    # plus
    textbox(s, Inches(6.35), Inches(3.6), Inches(0.5), Inches(0.5), "+", size=32, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    # sticky
    sticky = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.7), Inches(5.8), Inches(4.8))
    fill_solid(sticky, LIGHT_CORAL)
    sticky.adjustments[0] = 0.08
    textbox(s, Inches(7.2), Inches(1.95), Inches(5), Inches(0.3), "STICKY NOTES", size=11, bold=True, color=CORAL)
    textbox(s, Inches(7.2), Inches(2.3), Inches(5), Inches(0.5), "Vector DB (Chroma)", size=24, bold=True, color=NAVY)
    note = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(3.3), Inches(4.4), Inches(1.8))
    fill_solid(note, RGBColor(0xFF, 0xED, 0xD5))
    textbox(s, Inches(7.75), Inches(3.7), Inches(4), Inches(1.2), "“…Cap’n Bill and the Tin Woodman walked…”\n\nPassage match by meaning, not keywords", size=13, color=CORAL)
    textbox(s, Inches(7.2), Inches(5.7), Inches(5.2), Inches(0.5), "Exact wording when the graph path is fuzzy", size=13, color=CORAL)

    # Hybrid flow — plain English steps
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    textbox(s, Inches(0.55), Inches(0.28), Inches(12), Inches(0.28), "05  ·  HYBRID ASK FLOW", size=11, bold=True, color=TEAL2)
    textbox(s, Inches(0.55), Inches(0.58), Inches(12), Inches(0.5), "How an answer is born", size=26, bold=True, color=WHITE)
    textbox(
        s,
        Inches(0.55),
        Inches(1.1),
        Inches(12),
        Inches(0.35),
        "Corkboard first when the trail is strong — sticky notes when it’s not — never invent.",
        size=14,
        color=RGBColor(0xB8, 0xCB, 0xD1),
    )
    ask_steps = [
        ("1. Question", "Someone asks", TEAL, "A real question\nfrom a user."),
        ("2. Find names", "Who / what", TEAL2, "Spot people,\nproducts, places."),
        ("3. Corkboard", "Graph path", TEAL, "Follow links on\nthe knowledge map."),
        ("4. Sticky notes", "Passages", TEAL2, "Pull matching\nquotes for wording."),
        ("5. Merge", "Best evidence", CORAL, "Rank what actually\nsupports an answer."),
        ("6. Proof", "Trust Trail", CORAL, "Show sources —\nor refuse."),
    ]
    flow_row(s, ask_steps, top=Inches(1.65), height=Inches(3.35))
    textbox(
        s,
        Inches(0.55),
        Inches(5.3),
        Inches(12.2),
        Inches(1.4),
        "Strong trail → corkboard leads\n"
        "Weak trail → sticky notes / top passages lead\n"
        "No evidence → clarify or refuse (don’t invent)",
        size=15,
        color=WHITE,
    )

    # Why better
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_chrome(s)
    content_title(s, "06  ·  WHY SMARTER THAN VECTOR DB ALONE", "Vector is necessary. Hybrid is smarter.")
    card(
        s,
        Inches(0.55),
        Inches(1.55),
        Inches(4.0),
        Inches(4.9),
        "Vector only",
        "Good at “pages that sound alike.”\n\n"
        "Weak at who connects to whom.\n"
        "Hard to prove multi-step answers.\n"
        "Conflicts get smoothed over.\n"
        "Audit trail = a similarity score.",
        fill=LIGHT_CORAL,
        title_color=CORAL,
        body_size=13,
    )
    card(
        s,
        Inches(4.8),
        Inches(1.55),
        Inches(7.8),
        Inches(4.9),
        "VERA Hybrid = map + wording + proof",
        "• Relationships are on the corkboard — not guessed\n"
        "• Trust Trail managers and clients can follow\n"
        "• Conflicts listed with both sides + sources\n"
        "• Same sticky-note power when the map is thin\n"
        "• Smarter ≠ bigger model — smarter = structure + evidence + discipline\n"
        "• This is the industry direction (GraphRAG + vector together)",
        fill=LIGHT_TEAL,
        title_color=TEAL,
        body_size=13,
    )

    # Early Studio screens
    for kicker, title, img, blurb in [
        (
            "07  ·  STUDIO · HOME",
            "The trust dashboard",
            "01-home-fleet.png",
            "See how ready each agent is: grounded answers, evidence coverage, conflicts, and AI Findings you can Prove it on.",
        ),
        (
            "08  ·  PROVE IT",
            "Receipts, not slogans",
            "02-prove-it-drawer.png",
            "Click a finding → open the real list. Names, concepts, evidence edges with quotes, or conflicts with both sides.",
        ),
        (
            "09  ·  CONNECT",
            "Just bring your data",
            "04-connect.png",
            "Upload files, crawl a website, or connect SharePoint / Blob. Same pipeline after that: clean → weave → embed.",
        ),
    ]:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_light_chrome(s)
        content_title(s, kicker, title)
        add_image_if(s, img, Inches(0.45), Inches(1.35), Inches(8.6), max_height=Inches(5.5))
        card(s, Inches(9.2), Inches(1.35), Inches(3.65), Inches(5.5), "In plain English", blurb, body_size=13)

    # 10 · Librarian story (1–2 liners) → 10B · full-bleed Knowledge Map
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    textbox(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.35), "10  ·  THE LIBRARIAN", size=14, bold=True, color=TEAL2)
    textbox(
        s,
        Inches(0.7),
        Inches(1.9),
        Inches(11.8),
        Inches(1.5),
        "Most AI skims the shelves\nand answers with confidence.",
        size=30,
        bold=True,
        color=WHITE,
    )
    textbox(
        s,
        Inches(0.7),
        Inches(3.55),
        Inches(11.8),
        Inches(0.7),
        "VERA is the careful librarian.",
        size=32,
        bold=True,
        color=CORAL,
    )
    textbox(
        s,
        Inches(0.7),
        Inches(4.5),
        Inches(11.5),
        Inches(1.2),
        "Pins every name on a corkboard. Ties a string only when the books actually connect them. "
        "Keeps sticky-note quotes as receipts — so Cap’n Bill isn’t a guess; it’s a path you can follow.",
        size=17,
        color=RGBColor(0xC8, 0xD8, 0xDE),
    )
    textbox(
        s,
        Inches(0.7),
        Inches(6.0),
        Inches(11.5),
        Inches(0.5),
        "Next → see the corkboard VERA built from the Oz novels.",
        size=16,
        bold=True,
        color=TEAL2,
    )

    # Full-bleed map hero — maximize graph for presentation clarity
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, RGBColor(0xF8, 0xF9, 0xFB))
    band = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, Inches(0.72))
    fill_solid(band, DARK)
    textbox(s, Inches(0.4), Inches(0.12), Inches(9), Inches(0.28), "11  ·  KNOWLEDGE MAP", size=12, bold=True, color=TEAL2)
    textbox(
        s,
        Inches(0.4),
        Inches(0.36),
        Inches(12.5),
        Inches(0.3),
        "Frank Baum – Novel  ·  the librarian’s corkboard made visible",
        size=16,
        bold=True,
        color=WHITE,
    )
    hero = "06-knowledge-map-hero.png"
    if not (SHOTS / hero).exists():
        hero = "06-knowledge-map.png"
    add_image_if(
        s,
        hero,
        Inches(0.15),
        Inches(0.8),
        Inches(13.0),
        max_height=Inches(6.35),
    )
    textbox(
        s,
        Inches(0.4),
        Inches(7.15),
        Inches(12.5),
        Inches(0.28),
        "Colors = entity types  ·  lines = relationships  ·  denser clusters = richer connections in the books",
        size=11,
        color=MUTED,
    )

    # Remaining Studio screens
    for kicker, title, img, blurb in [
        (
            "12  ·  ASK",
            "Daily Q&A with proof",
            "05-ask.png",
            "Ask a question. VERA answers from the corkboard + sticky notes, then shows the Trust Trail (or refuses if evidence is thin).",
        ),
        (
            "13  ·  FLEET",
            "Many specialists, one studio",
            "03-fleet.png",
            "Create agents for different topics. Each has its own knowledge. Track Draft → Ready → Live before you publish.",
        ),
        (
            "14  ·  DEPLOY & EMBED",
            "Put VERA on any website",
            "09-deploy.png",
            "Publish → copy a small HTML / iframe snippet → paste on a marketing site or intranet. Same trusted brain.",
        ),
    ]:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_light_chrome(s)
        content_title(s, kicker, title)
        add_image_if(s, img, Inches(0.45), Inches(1.35), Inches(8.6), max_height=Inches(5.5))
        card(s, Inches(9.2), Inches(1.35), Inches(3.65), Inches(5.5), "In plain English", blurb, body_size=13)

    # Embed architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    textbox(s, Inches(0.55), Inches(0.3), Inches(12), Inches(0.28), "15  ·  EMBED ARCHITECTURE", size=11, bold=True, color=TEAL2)
    textbox(s, Inches(0.55), Inches(0.65), Inches(12), Inches(0.55), "Website → VERA → your knowledge", size=26, bold=True, color=WHITE)
    textbox(
        s,
        Inches(0.55),
        Inches(1.2),
        Inches(12),
        Inches(0.35),
        "One agent brain. Many places customers can ask it.",
        size=15,
        color=RGBColor(0xB8, 0xCB, 0xD1),
    )
    embed_steps = [
        (
            "Client website",
            "Chat widget",
            CORAL,
            "Visitor asks on your site\n(HTML, iframe, or API).\nLooks like your product.",
        ),
        (
            "VERA agent",
            "The librarian",
            TEAL,
            "Hybrid Ask + Trust Trail.\nSecure embed key.\nAnswers only from this agent.",
        ),
        (
            "Your knowledge",
            "Private library",
            TEAL2,
            "Docs, site, SharePoint.\nIsolated per agent.\nUpload once → serve everywhere.",
        ),
    ]
    flow_row(s, embed_steps, top=Inches(1.75), height=Inches(3.25), gap=Inches(0.4))
    textbox(
        s,
        Inches(0.55),
        Inches(5.3),
        Inches(12.2),
        Inches(1.2),
        "Publish in Studio → copy snippet → paste on any site.\nUpload once → chatbot wherever your clients work.",
        size=16,
        color=WHITE,
    )

    # Feature map — short “what it’s for” on every tile
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_chrome(s)
    content_title(s, "16  ·  FEATURE MAP", "Everything in one glance — what each piece is for")
    feats = [
        ("Agents", "A specialist AI for one topic. Its knowledge stays private to that agent."),
        ("Connect", "Way in for your data: upload files, crawl a site, SharePoint, or Blob."),
        ("CleanStack", "Cleans messy docs first so the librarian isn’t reading junk."),
        ("Graph", "The corkboard — people, products, places, and how they connect."),
        ("Chroma", "Sticky-note quotes — exact wording stored for meaning search."),
        ("Hybrid Ask", "Uses corkboard + sticky notes together for grounded answers."),
        ("Trust Trail", "The receipt: hops, quotes, and sources behind the answer."),
        ("Prove it", "Click any finding to open the real evidence list underneath."),
        ("Conflicts", "When sources disagree — show both sides, don’t paper over it."),
        ("Map", "See the corkboard visually — what VERA learned from your books."),
        ("Publish", "Embed the same agent on any website as a trusted chat widget."),
        ("Refuse", "If evidence is missing, say so. No empty confidence."),
    ]
    for i, (t, b) in enumerate(feats):
        r, c = divmod(i, 4)
        accent = CORAL if t in ("Prove it", "Refuse", "Conflicts") else TEAL
        card(
            s,
            Inches(0.4 + c * 3.2),
            Inches(1.35 + r * 1.85),
            Inches(3.1),
            Inches(1.75),
            t,
            b,
            title_color=accent,
            body_size=11,
        )

    # Results scorecard (3 corpora) — matches manager-edited numbers
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_chrome(s)
    content_title(s, "17  ·  MEASURED RESULTS", "Real corpora — not vanity demos")
    card(
        s,
        Inches(0.45),
        Inches(1.45),
        Inches(4.05),
        Inches(5.3),
        "PlayReady · ~88%*",
        "Product/docs KB (41 PDFs).\n\n"
        "133/150 on PublicBot golden IDs\n"
        "(licensing, EV cert, distribution).\n\n"
        "≈88% on the evaluated set.\n\n"
        "*Directional — keep expanding hard cases.",
        fill=LIGHT_TEAL,
        title_color=TEAL,
        body_size=12,
    )
    card(
        s,
        Inches(4.65),
        Inches(1.45),
        Inches(4.05),
        Inches(5.3),
        "Oz novels · ~90%",
        "Document/narrative KB (12 full-text books).\n\n"
        "181/200 on oz_baum_v1 golden suite (≈90%).\n\n"
        "Strong on character & relationship questions\n"
        "where graph + narrative help.\n\n"
        "Gaps: hard multi-character recall.",
        fill=WHITE,
        title_color=NAVY,
        body_size=12,
    )
    card(
        s,
        Inches(8.85),
        Inches(1.45),
        Inches(4.0),
        Inches(5.3),
        "Thoughtworks · ~56%",
        "Public website crawl (~442 pages).\n\n"
        "~55–56% on thoughtworks web golden\n"
        "(leaders, services, hard FAQ).\n\n"
        "Hardest corpus: web chrome, thin pages,\n"
        "multi-intent questions.\n\n"
        "Honesty: trust still matters when % is mid.",
        fill=LIGHT_CORAL,
        title_color=CORAL,
        body_size=12,
    )

    # Future
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_chrome(s, CORAL)
    content_title(s, "18  ·  FUTURE USE CASES", "Day-in, day-out client value")
    futures = [
        ("Trusted site assistant", "Staff and customers ask on your site — answers come with citations, not vibes."),
        ("“Is this true?” checker", "Paste a claim → Supported / Conflicted / Unknown, with sources."),
        ("Reply helper", "Draft support or sales replies from approved knowledge only."),
        ("Compliance coach", "See obligations and conflicts clearly for regulated teams."),
        ("Onboarding paths", "Role-based “what must I know?” journeys from your library."),
        ("More connectors", "Outlook, Teams, Drive — same ingest pattern, same trust layer."),
    ]
    for i, (t, b) in enumerate(futures):
        r, c = divmod(i, 3)
        card(
            s,
            Inches(0.55 + c * 4.2),
            Inches(1.55 + r * 2.5),
            Inches(4.0),
            Inches(2.3),
            t,
            b,
            body_size=13,
        )

    # Client story
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    textbox(s, Inches(0.55), Inches(0.35), Inches(12), Inches(0.28), "19  ·  CLIENT STORY", size=11, bold=True, color=TEAL2)
    textbox(s, Inches(0.55), Inches(0.75), Inches(12), Inches(0.55), "Upload → magic → trusted chatbot", size=28, bold=True, color=WHITE)
    story = [
        ("1. Bring data", "Your library", TEAL, "Upload files, crawl a site,\nor connect SharePoint."),
        ("2. See magic", "Corkboard appears", TEAL2, "Graph, Map, Prove it —\nwatch the librarian organize."),
        ("3. Ask + proof", "Daily Q&A", CORAL, "Answers with Trust Trail.\nRefuse when unsure."),
        ("4. Embed", "On their site", RGBColor(0xF4, 0x8C, 0x06), "Paste a widget —\nsame specialist for customers."),
    ]
    flow_row(s, story, top=Inches(1.7), height=Inches(3.4), gap=Inches(0.35))
    textbox(s, Inches(0.55), Inches(5.5), Inches(12), Inches(0.8), "Just bring your data. VERA does the rest.", size=22, bold=True, color=WHITE)

    # Close
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    textbox(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.4), "CLOSING", size=12, bold=True, color=TEAL2)
    textbox(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(1.3), "Evidence you can live with — every day", size=34, bold=True, color=WHITE)
    textbox(
        s,
        Inches(0.8),
        Inches(3.8),
        Inches(11.2),
        Inches(1.4),
        "Not a louder chatbot.\n"
        "A careful librarian: corkboard + sticky notes + receipts —\n"
        "ready to embed on any website.",
        size=18,
        color=RGBColor(0xC8, 0xD8, 0xDE),
    )
    textbox(s, Inches(0.8), Inches(5.6), Inches(11), Inches(0.5), "Live demo path: Home → Prove it → Map → Ask → Publish", size=16, bold=True, color=CORAL)

    from datetime import datetime

    stamped = ROOT / f"VERA-Complete-{datetime.now().strftime('%Y%m%d-%H%M%S')}.pptx"
    target = None
    for path in [OUT, *OUT_FALLBACKS, stamped]:
        try:
            prs.save(path)
            target = path
            if path != OUT:
                print(f"Note: earlier PPTX file(s) locked — wrote {path.name}")
            break
        except PermissionError:
            continue
    if target is None:
        raise PermissionError("Close open PowerPoint files under docs/presentations/ and retry.")
    print(f"Wrote {target} ({target.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    build()
