"""PDF / Office / text / image OCR extractors (documents only — no HTML crawl)."""

from __future__ import annotations

import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

_MIN_ALPHA_CHARS = 40
_OCR_MAX_PAGES = 30
_OCR_ZOOM = 2.0
_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp"}


def extract_document_text(filename: str, mime: str | None, content: bytes) -> tuple[str, dict]:
    """Extract text from an uploaded/sharepoint/blob document. No HTML path."""
    suffix = Path(filename).suffix.lower()
    mime = mime or ""

    if suffix in {".txt", ".md"} or mime.startswith("text/"):
        return content.decode("utf-8", errors="ignore"), {"kind": "text"}

    if suffix == ".pdf" or mime == "application/pdf":
        text, structure = _extract_pdf_text(content)
        if _needs_ocr(text, int(structure.get("pages") or 0)):
            ocr_text, ocr_meta = _ocr_pdf(content)
            if _alpha_count(ocr_text) > _alpha_count(text):
                return ocr_text, {**structure, **ocr_meta, "ocr": True}
            if text.strip():
                structure["ocr_attempted"] = True
                return text, structure
            return ocr_text, {**structure, **ocr_meta, "ocr": True}
        return text, structure

    if suffix in _IMAGE_SUFFIXES or mime.startswith("image/"):
        return _ocr_image_bytes(content)

    if suffix == ".docx" or mime.endswith("wordprocessingml.document"):
        import docx

        document = docx.Document(io.BytesIO(content))
        paras = [p.text for p in document.paragraphs if p.text.strip()]
        return "\n".join(paras), {"kind": "docx", "paragraphs": len(paras)}

    if suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        slides = []
        for i, slide in enumerate(prs.slides):
            bits = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    bits.append(shape.text)
            slides.append({"slide": i + 1, "text": "\n".join(bits)})
        return "\n\n".join(s["text"] for s in slides), {"kind": "pptx", "slides": len(slides)}

    if suffix == ".xlsx":
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                vals = [str(c) for c in row if c is not None]
                if vals:
                    rows.append(" | ".join(vals))
            parts.append(f"# Sheet {sheet.title}\n" + "\n".join(rows[:200]))
        return "\n\n".join(parts), {"kind": "xlsx", "sheets": len(wb.worksheets)}

    return content.decode("utf-8", errors="ignore"), {"kind": "unknown"}


def needs_ocr_emit(filename: str, mime: str | None, content: bytes) -> bool:
    """True when PDF/image path will run OCR (for progress events)."""
    suffix = Path(filename).suffix.lower()
    mime = mime or ""
    if suffix in _IMAGE_SUFFIXES or mime.startswith("image/"):
        return True
    if suffix == ".pdf" or mime == "application/pdf":
        text, structure = _extract_pdf_text(content)
        return _needs_ocr(text, int(structure.get("pages") or 0))
    return False


def _extract_pdf_text(content: bytes) -> tuple[str, dict]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages = []
    for i, page in enumerate(reader.pages):
        pages.append({"page": i + 1, "text": page.extract_text() or ""})
    text = "\n\n".join(p["text"] for p in pages)
    return text, {"kind": "pdf", "pages": len(pages)}


def _needs_ocr(text: str, page_count: int) -> bool:
    if page_count <= 0:
        return True
    if _alpha_count(text) < _MIN_ALPHA_CHARS:
        return True
    if page_count >= 2 and _alpha_count(text) < page_count * 15:
        return True
    return False


def _alpha_count(text: str) -> int:
    return sum(1 for c in text if c.isalpha())


def _ocr_pdf(content: bytes) -> tuple[str, dict]:
    import fitz  # PyMuPDF
    import pytesseract
    from PIL import Image

    doc = fitz.open(stream=content, filetype="pdf")
    try:
        n = min(doc.page_count, _OCR_MAX_PAGES)
        parts: list[str] = []
        for i in range(n):
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=fitz.Matrix(_OCR_ZOOM, _OCR_ZOOM), alpha=False)
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            page_text = pytesseract.image_to_string(img) or ""
            if page_text.strip():
                parts.append(f"[Page {i + 1}]\n{page_text.strip()}")
        text = "\n\n".join(parts)
        return text, {
            "kind": "pdf",
            "pages": doc.page_count,
            "ocr_pages": n,
            "ocr_engine": "tesseract",
        }
    finally:
        doc.close()


def _ocr_image_bytes(content: bytes) -> tuple[str, dict]:
    import pytesseract
    from PIL import Image

    img = Image.open(io.BytesIO(content))
    if img.mode not in {"RGB", "L"}:
        img = img.convert("RGB")
    text = pytesseract.image_to_string(img) or ""
    return text.strip(), {
        "kind": "image",
        "ocr": True,
        "ocr_pages": 1,
        "ocr_engine": "tesseract",
        "size": list(img.size),
    }
