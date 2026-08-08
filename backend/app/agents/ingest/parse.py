from __future__ import annotations

import hashlib
import io
import logging
from pathlib import Path

from app.agents.base import AgentContext, AgentError, AgentResult, AgentWarning
from app.agents.ingest.contracts import ParseInput, ParseOutput, ParsedFile

logger = logging.getLogger(__name__)

# Sparse digital text → treat as scanned and OCR
_MIN_ALPHA_CHARS = 40
_OCR_MAX_PAGES = 30
_OCR_ZOOM = 2.0
_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp"}


class ParseAgent:
    id = "parse"
    display_name = "Parse Agent"
    input_model = ParseInput
    output_model = ParseOutput

    async def run(self, ctx: AgentContext, payload: ParseInput) -> AgentResult[ParseOutput]:
        store = ctx.stores
        parsed: list[ParsedFile] = []
        warnings: list[AgentWarning] = []

        for f in payload.files:
            try:
                ctx.emit(self.id, "parse.file", f"Parsing {f.filename}", progress=0.2)
                await ctx.flush_job_progress(0.2)
                text, structure = await _extract_with_ocr(ctx, f.filename, f.mime, f.content)
            except Exception as exc:  # noqa: BLE001
                logger.warning("Parse failed for %s: %s", f.filename, exc)
                warnings.append(
                    AgentWarning(code="PARSE_FAIL", message=f"{f.filename}: {exc}")
                )
                text, structure = "", {"error": str(exc)}

            text = text.strip()
            text_hash = hashlib.sha256(text.encode("utf-8", errors="ignore")).hexdigest()
            if store is not None:
                await store.update_source_instance(
                    ctx.workspace_id,
                    f.source_id,
                    text_hash=text_hash,
                    status="parsed",
                )

            file_warnings: list[str] = []
            if not text:
                file_warnings.append("empty_extract")
            if structure.get("ocr"):
                file_warnings.append("ocr_used")
                warnings.append(
                    AgentWarning(
                        code="OCR_USED",
                        message=f"{f.filename}: used OCR ({structure.get('ocr_pages', 0)} pages)",
                    )
                )

            parsed.append(
                ParsedFile(
                    source_id=f.source_id,
                    filename=f.filename,
                    mime=f.mime,
                    text=text,
                    binary_hash=f.binary_hash,
                    text_hash=text_hash,
                    byte_size=f.byte_size,
                    appears_at=f.appears_at,
                    structure=structure,
                    warnings=file_warnings,
                )
            )

        if store is not None:
            await store.commit()

        nonempty = sum(1 for p in parsed if p.text)
        if nonempty == 0:
            return AgentResult(
                ok=False,
                error=AgentError(
                    code="PARSE_EMPTY",
                    message=(
                        "No extractable text even after OCR. "
                        "Check the file is a readable PDF/image (not blank or encrypted)."
                    ),
                ),
                warnings=warnings,
            )

        ctx.emit(self.id, "parse.done", f"Parsed {nonempty}/{len(parsed)} files", progress=1.0)
        await ctx.flush_job_progress(0.35)
        return AgentResult(
            ok=True,
            data=ParseOutput(files=parsed),
            warnings=warnings,
            metrics={"parsed": nonempty, "total": len(parsed)},
        )


async def _extract_with_ocr(
    ctx: AgentContext, filename: str, mime: str | None, content: bytes
) -> tuple[str, dict]:
    suffix = Path(filename).suffix.lower()
    mime = mime or ""

    if suffix in {".txt", ".md"} or mime.startswith("text/"):
        return content.decode("utf-8", errors="ignore"), {"kind": "text"}

    if suffix == ".pdf" or mime == "application/pdf":
        text, structure = _extract_pdf_text(content)
        if _needs_ocr(text, int(structure.get("pages") or 0)):
            ctx.emit(
                "parse",
                "parse.ocr",
                f"Scanned/image PDF detected — running OCR on {filename}",
                progress=0.28,
            )
            await ctx.flush_job_progress(0.28)
            ocr_text, ocr_meta = _ocr_pdf(content)
            if _alpha_count(ocr_text) > _alpha_count(text):
                return ocr_text, {**structure, **ocr_meta, "ocr": True}
            if text.strip():
                structure["ocr_attempted"] = True
                return text, structure
            return ocr_text, {**structure, **ocr_meta, "ocr": True}
        return text, structure

    if suffix in _IMAGE_SUFFIXES or mime.startswith("image/"):
        ctx.emit("parse", "parse.ocr", f"OCR image {filename}", progress=0.28)
        await ctx.flush_job_progress(0.28)
        return _ocr_image_bytes(content)

    if suffix == ".docx" or mime.endswith("wordprocessingml.document"):
        import docx

        document = docx.Document(io.BytesIO(content))
        paras = [p.text for p in document.paragraphs if p.text.strip()]
        return "\n".join(paras), {"kind": "docx", "paragraphs": len(paras)}

    if suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        slides = []
        for i, slide in enumerate(prs.slides):
            bits = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    bits.append(shape.text)
            slides.append({"slide": i + 1, "text": "\n".join(bits)})
        return "\n\n".join(s["text"] for s in slides), {"kind": "pptx", "slides": len(slides)}

    if suffix == ".xlsx":
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                vals = [str(c) for c in row if c is not None]
                if vals:
                    rows.append(" | ".join(vals))
            parts.append(f"# Sheet {sheet.title}\n" + "\n".join(rows[:200]))
        return "\n\n".join(parts), {"kind": "xlsx", "sheets": len(wb.worksheets)}

    return content.decode("utf-8", errors="ignore"), {"kind": "unknown"}


def _extract_pdf_text(content: bytes) -> tuple[str, dict]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages = []
    for i, page in enumerate(reader.pages):
        pages.append({"page": i + 1, "text": page.extract_text() or ""})
    text = "\n\n".join(p["text"] for p in pages)
    return text, {"kind": "pdf", "pages": len(pages)}


def _needs_ocr(text: str, page_count: int) -> bool:
    if page_count <= 0:
        return True
    if _alpha_count(text) < _MIN_ALPHA_CHARS:
        return True
    # Extremely sparse relative to page count (likely scanned)
    if page_count >= 2 and _alpha_count(text) < page_count * 15:
        return True
    return False


def _alpha_count(text: str) -> int:
    return sum(1 for c in text if c.isalpha())


def _ocr_pdf(content: bytes) -> tuple[str, dict]:
    import fitz  # PyMuPDF
    import pytesseract
    from PIL import Image

    doc = fitz.open(stream=content, filetype="pdf")
    try:
        n = min(doc.page_count, _OCR_MAX_PAGES)
        parts: list[str] = []
        for i in range(n):
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=fitz.Matrix(_OCR_ZOOM, _OCR_ZOOM), alpha=False)
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            page_text = pytesseract.image_to_string(img) or ""
            if page_text.strip():
                parts.append(f"[Page {i + 1}]\n{page_text.strip()}")
        text = "\n\n".join(parts)
        return text, {
            "kind": "pdf",
            "pages": doc.page_count,
            "ocr_pages": n,
            "ocr_engine": "tesseract",
        }
    finally:
        doc.close()


def _ocr_image_bytes(content: bytes) -> tuple[str, dict]:
    import pytesseract
    from PIL import Image

    img = Image.open(io.BytesIO(content))
    if img.mode not in {"RGB", "L"}:
        img = img.convert("RGB")
    text = pytesseract.image_to_string(img) or ""
    return text.strip(), {
        "kind": "image",
        "ocr": True,
        "ocr_pages": 1,
        "ocr_engine": "tesseract",
        "size": list(img.size),
    }
